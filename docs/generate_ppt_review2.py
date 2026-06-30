"""
Generate Review 2 PowerPoint Presentation for Smart Contract Auditor project.
Run: python generate_ppt_review2.py
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour Palette ───────────────────────────────────────────────
BG_DARK = RGBColor(0x0F, 0x0F, 0x23)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0xD4, 0xFF)
WHITE = RGBColor(0xE0, 0xE0, 0xE0)
GRAY = RGBColor(0xA0, 0xA0, 0xB0)
RED = RGBColor(0xFF, 0x47, 0x57)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
YELLOW = RGBColor(0xFF, 0xA5, 0x02)
GREEN = RGBColor(0x2E, 0xD5, 0x73)
BLUE = RGBColor(0x1E, 0x90, 0xFF)
PURPLE = RGBColor(0x7C, 0x5C, 0xFF)


# ── Helper Functions ─────────────────────────────────────────────

def set_slide_bg(slide, color=BG_DARK):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text,
                 font_size=18, color=WHITE, bold=False,
                 alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_section_title(slide, title):
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 title, font_size=34, color=ACCENT, bold=True)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(3), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_bullet_slide(slide, items, start_top=Inches(1.8),
                     font_size=16, color=WHITE, bullet_color=ACCENT):
    for i, item in enumerate(items):
        top = start_top + Inches(i * 0.5)
        txBox = slide.shapes.add_textbox(Inches(1), top, Inches(10.5), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        run1 = p.add_run()
        run1.text = "  \u25b8  "
        run1.font.size = Pt(font_size)
        run1.font.color.rgb = bullet_color
        run1.font.bold = True

        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = color


def add_card(slide, left, top, width, height, title, body, border_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(title.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14) if idx == 0 else Pt(11)
        p.font.color.rgb = border_color if idx == 0 else GRAY
        p.font.bold = (idx == 0)
        p.alignment = PP_ALIGN.CENTER

    if body:
        for line in body.split("\n"):
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(10)
            p.font.color.rgb = GRAY
            p.alignment = PP_ALIGN.CENTER


# ═════════════════════════════════════════════════════════════════
#                     PRESENTATION BUILDER
# ═════════════════════════════════════════════════════════════════

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 1 — Title
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
                 "Smart Contract Auditor",
                 font_size=46, color=ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.8),
                 "AI-Powered Vulnerability Detection for Solidity Smart Contracts",
                 font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.5),
                 "Leveraging Static Analysis \u2022 Control-Flow Graphs \u2022 Taint Propagation \u2022 Machine Learning",
                 font_size=17, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.5),
                 "Software Design Project \u2014 Review 2 (Final)",
                 font_size=20, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                 "Student: Harish  |  Guide: [Your Guide Name]",
                 font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.5),
                 "Department of Computer Science and Engineering",
                 font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 2 — Recap: Problem & Motivation
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Problem & Motivation")

    items = [
        "DeFi ecosystems hold billions in Total Value Locked (TVL) \u2014 a single smart contract bug can drain entire protocols",
        "$3.8 Billion lost to blockchain exploits in 2022 alone (Chainalysis Report)",
        "Traditional audit tools like Slither and Mythril suffer from high false-positive rates",
        "Manual audits cost $5K\u2013$50K+ and take weeks \u2014 infeasible for most developers",
        "Gap: No open tool combines hybrid analysis, ML-based ranking, and human-readable explanations",
        "Our Goal: Build an end-to-end AI auditor that is accurate, explainable, and accessible",
    ]
    add_bullet_slide(slide, items, start_top=Inches(1.8), font_size=15)

    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
                 "\"Security must be automated, explainable, and affordable for Web3 to scale.\"",
                 font_size=16, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 3 — Review 1 Recap & What Changed
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Review 1 Recap \u2192 What\u2019s New in Review 2")

    # Review 1 column
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
                 "Review 1 (Completed)", font_size=18, color=GREEN, bold=True)

    r1_items = [
        "15-rule Static Analyzer (regex pattern matching)",
        "ML Classifier with heuristic-based scoring",
        "Explainability Engine (attack scenarios, fix examples)",
        "FastAPI REST API + React frontend",
        "CLI tool for batch scanning",
        "HTML & JSON report generation",
    ]
    for i, item in enumerate(r1_items):
        top = Inches(2.2) + Inches(i * 0.42)
        txBox = slide.shapes.add_textbox(Inches(1), top, Inches(5.2), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = "\u2713  "
        run1.font.size = Pt(14)
        run1.font.color.rgb = GREEN
        run1.font.bold = True
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(13)
        run2.font.color.rgb = GRAY

    # Review 2 column
    add_text_box(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(0.5),
                 "Review 2 (New)", font_size=18, color=ACCENT, bold=True)

    r2_items = [
        "Control-Flow Graph (CFG) Analyzer",
        "Inter-procedural Taint Propagation Engine",
        "Gas Optimization Analyzer",
        "Comprehensive Test Suite (8 test modules)",
        "Edge-case & regression testing library",
        "Docker containerisation (multi-service)",
    ]
    for i, item in enumerate(r2_items):
        top = Inches(2.2) + Inches(i * 0.42)
        txBox = slide.shapes.add_textbox(Inches(7), top, Inches(5.2), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = "\u2605  "
        run1.font.size = Pt(14)
        run1.font.color.rgb = ACCENT
        run1.font.bold = True
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(13)
        run2.font.color.rgb = WHITE

    add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.5),
                 "Analysis depth increased from single-pass regex to multi-layer graph-based + data-flow detection",
                 font_size=14, color=YELLOW, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 4 — System Architecture (Final)
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Final System Architecture")

    # Row 1 — Interfaces
    add_text_box(slide, Inches(0.5), Inches(1.6), Inches(12), Inches(0.4),
                 "User Interfaces", font_size=13, color=GRAY, bold=True)
    add_card(slide, Inches(0.5), Inches(2.0), Inches(3.7), Inches(1.1),
             "React Frontend\n(Dark-theme Dashboard)", None, BLUE)
    add_card(slide, Inches(4.5), Inches(2.0), Inches(3.7), Inches(1.1),
             "FastAPI REST API\n(/analyze, /batch, /report)", None, PURPLE)
    add_card(slide, Inches(8.7), Inches(2.0), Inches(3.7), Inches(1.1),
             "CLI Tool\n(Single & Batch Scanning)", None, GRAY)

    # Row 2 — Core Analysis Pipeline
    add_text_box(slide, Inches(0.5), Inches(3.3), Inches(12), Inches(0.4),
                 "Core Analysis Pipeline", font_size=13, color=GRAY, bold=True)
    add_card(slide, Inches(0.5), Inches(3.7), Inches(2.6), Inches(1.3),
             "Static Analyzer\n(15 SWC Rules)", "Regex-based\npattern detection", ORANGE)
    add_card(slide, Inches(3.3), Inches(3.7), Inches(2.6), Inches(1.3),
             "CFG Analyzer\n(Control-Flow Graphs)", "Path-sensitive\nvulnerability tracking", GREEN)
    add_card(slide, Inches(6.1), Inches(3.7), Inches(2.6), Inches(1.3),
             "Taint Analyzer\n(Data-Flow Tracking)", "Source-to-Sink\npropagation analysis", RED)
    add_card(slide, Inches(8.9), Inches(3.7), Inches(3.5), Inches(1.3),
             "ML Classifier\n(Risk Scoring)", "20-feature extraction\nhybrid confidence", YELLOW)

    # Row 3 — Output
    add_text_box(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(0.4),
                 "Post-Processing & Output", font_size=13, color=GRAY, bold=True)
    add_card(slide, Inches(0.5), Inches(5.6), Inches(3.0), Inches(1.0),
             "Explainability Engine", "Attack narratives + fix examples", ACCENT)
    add_card(slide, Inches(3.8), Inches(5.6), Inches(3.0), Inches(1.0),
             "Report Generator", "HTML \u2022 JSON \u2022 Text", PURPLE)
    add_card(slide, Inches(7.1), Inches(5.6), Inches(2.8), Inches(1.0),
             "Gas Optimizer", "Efficiency scoring", BLUE)
    add_card(slide, Inches(10.2), Inches(5.6), Inches(2.3), Inches(1.0),
             "Analysis Store", "In-memory / DB", GRAY)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 5 — CFG Analyzer Deep Dive
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "NEW: Control-Flow Graph (CFG) Analyzer")

    steps = [
        ("What is a CFG?",
         "A directed graph where each node represents a basic block of code "
         "(assignment, branch, loop, external call) and edges represent possible "
         "execution paths. This enables path-sensitive vulnerability detection that "
         "regex pattern matching alone cannot achieve."),
        ("How We Build It",
         "1.  Extract every function from the contract source\n"
         "2.  Tokenise each function body into basic blocks (entry, branch, loop, call, exit)\n"
         "3.  Link blocks via successor / predecessor edges based on control-flow semantics\n"
         "4.  Annotate nodes with properties: external_call, state_change, ether_transfer, access_control"),
        ("Path-Sensitive Detection",
         "\u2022 Reentrancy: external call node reachable before state-change node on the same path\n"
         "\u2022 Unguarded Ether Transfer: transfer node without an access-control predecessor\n"
         "\u2022 Cyclomatic Complexity: high branching factor \u2192 higher bug likelihood\n"
         "\u2022 Unreachable Code: nodes with no predecessors after entry"),
    ]

    for i, (title, desc) in enumerate(steps):
        top = Inches(1.7) + Inches(i * 1.7)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.5), Inches(1.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = GREEN
        shape.line.width = Pt(1)
        add_text_box(slide, Inches(1.2), top + Inches(0.08), Inches(10.5), Inches(0.35),
                     title, font_size=16, color=GREEN, bold=True)
        add_text_box(slide, Inches(1.2), top + Inches(0.4), Inches(10.5), Inches(0.95),
                     desc, font_size=12, color=GRAY)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 6 — Taint Analysis Deep Dive
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "NEW: Inter-Procedural Taint Propagation Engine")

    steps = [
        ("Concept: Source \u2192 Propagation \u2192 Sink",
         "Taint analysis tracks how untrusted (user-controlled) data flows through the contract. "
         "If attacker-controlled data reaches a dangerous operation without sanitisation, "
         "that data path constitutes a vulnerability."),
        ("Sources (Untrusted Inputs)",
         "msg.sender \u2022 msg.value \u2022 msg.data \u2022 tx.origin \u2022 function parameters \u2022 calldataload \u2022 external call return values\n"
         "Each source is tagged with a provenance label for traceability."),
        ("Sinks (Dangerous Operations)",
         "Ether transfer (.transfer / .send / .call{value}) \u2022 delegatecall target \u2022 selfdestruct beneficiary\n"
         "Unchecked storage writes \u2022 Array index without bounds check \u2022 Attacker-controlled guards"),
        ("Sanitisers & False-Positive Reduction",
         "require(), assert(), if-checks, and revert() are recognised as sanitisers. "
         "If a sanitiser sits between a source and a sink on the data-flow path, the finding "
         "is suppressed \u2014 significantly reducing false positives compared to pattern-only tools."),
    ]

    for i, (title, desc) in enumerate(steps):
        top = Inches(1.65) + Inches(i * 1.35)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.5), Inches(1.15))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = RED
        shape.line.width = Pt(1)
        add_text_box(slide, Inches(1.2), top + Inches(0.06), Inches(10.5), Inches(0.32),
                     title, font_size=15, color=RED, bold=True)
        add_text_box(slide, Inches(1.2), top + Inches(0.36), Inches(10.5), Inches(0.75),
                     desc, font_size=12, color=GRAY)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 7 — Hybrid Analysis Pipeline
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Hybrid Analysis Pipeline \u2014 How It All Fits Together")

    pipeline = [
        ("1", "Ingest", "Solidity source code uploaded via Web UI, API, or CLI", BLUE),
        ("2", "Static Analysis", "15 SWC rule-based pattern detectors scan for known vulnerability signatures", ORANGE),
        ("3", "CFG Construction", "Build control-flow graphs per function; annotate nodes with security properties", GREEN),
        ("4", "Taint Propagation", "Track untrusted inputs from sources through assignments to dangerous sinks", RED),
        ("5", "ML Classification", "Extract 20 code features \u2192 predict across 10 vulnerability classes", YELLOW),
        ("6", "Merge & Rank", "Combine results: 60% static + 40% ML confidence; deduplicate; sort by severity", PURPLE),
        ("7", "Explain & Report", "Generate attack narratives, fix examples, SWC references; export HTML/JSON", ACCENT),
    ]

    for i, (num, title, desc, color) in enumerate(pipeline):
        top = Inches(1.7) + Inches(i * 0.72)
        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), top + Inches(0.05), Inches(0.45), Inches(0.45))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        ctf = circle.text_frame
        cp = ctf.paragraphs[0]
        cp.text = num
        cp.font.size = Pt(14)
        cp.font.color.rgb = BG_DARK
        cp.font.bold = True
        cp.alignment = PP_ALIGN.CENTER

        add_text_box(slide, Inches(1.5), top, Inches(2), Inches(0.5),
                     title, font_size=15, color=color, bold=True)
        add_text_box(slide, Inches(3.8), top, Inches(8.5), Inches(0.5),
                     desc, font_size=13, color=WHITE)

    add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.5),
                 "Multi-layer approach catches vulnerabilities that any single technique would miss",
                 font_size=14, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 8 — ML Classification Details
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "ML Feature Extraction & Vulnerability Classification")

    # Left: Features
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
                 "20 Features Extracted from Solidity Code", font_size=16, color=ACCENT, bold=True)

    feature_groups = [
        ("Code Complexity", "LOC, function count, modifier count, event count, mapping count", BLUE),
        ("Control Flow", "require() count, assert() count, loop count", GREEN),
        ("External Interactions", "external call count, transfer count, send count, delegatecall count", RED),
        ("Security Patterns", "ReentrancyGuard, Ownable, SafeMath usage, tx.origin, selfdestruct", ORANGE),
        ("Solidity Metadata", "Version safety flag, floating pragma flag, state variable count", PURPLE),
    ]
    for i, (group, features, color) in enumerate(feature_groups):
        top = Inches(2.2) + Inches(i * 0.75)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(5.8), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = color
        shape.line.width = Pt(1)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{group}: "
        r.font.size = Pt(12)
        r.font.color.rgb = color
        r.font.bold = True
        r2 = p.add_run()
        r2.text = features
        r2.font.size = Pt(11)
        r2.font.color.rgb = GRAY

    # Right: Scoring
    add_text_box(slide, Inches(7.2), Inches(1.6), Inches(5.5), Inches(0.5),
                 "Hybrid Confidence Scoring", font_size=16, color=ACCENT, bold=True)

    scoring_items = [
        "10 vulnerability classes predicted simultaneously",
        "Heuristic scoring calibrated by domain experts",
        "Combined: 60% Static + 40% ML confidence",
        "ML-only findings added when confidence > 0.6",
        "Severity dynamically adjusted by final score",
        "Results sorted by descending confidence",
    ]
    for i, item in enumerate(scoring_items):
        top = Inches(2.3) + Inches(i * 0.55)
        txBox = slide.shapes.add_textbox(Inches(7.2), top, Inches(5.5), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = "\u25b8  "
        run1.font.size = Pt(13)
        run1.font.color.rgb = YELLOW
        run1.font.bold = True
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(13)
        run2.font.color.rgb = WHITE

    add_text_box(slide, Inches(0.5), Inches(6.3), Inches(12), Inches(0.8),
                 "Severity Thresholds:  Critical \u2265 0.85  |  High \u2265 0.70  |  Medium \u2265 0.50  |  Low \u2265 0.30  |  Info < 0.30",
                 font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 9 — Vulnerability Coverage
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Comprehensive Vulnerability Coverage")

    vulns = [
        ("SWC-107", "Reentrancy (CEI Violation)", "Critical", RED, "Static + CFG + Taint"),
        ("SWC-101", "Integer Overflow / Underflow", "High", ORANGE, "Static + ML"),
        ("SWC-115", "tx.origin Authorization", "High", ORANGE, "Static + Taint"),
        ("SWC-104", "Unchecked Call Return Value", "High", ORANGE, "Static"),
        ("SWC-105", "Unprotected Ether Withdrawal", "Critical", RED, "Static + CFG"),
        ("SWC-106", "Unprotected SELFDESTRUCT", "Critical", RED, "Static + Taint"),
        ("SWC-112", "Delegatecall to Untrusted Callee", "Critical", RED, "Static + Taint"),
        ("SWC-120", "Weak Randomness (block variables)", "High", ORANGE, "Static"),
        ("SWC-103", "Floating Pragma", "Info", BLUE, "Static"),
        ("SWC-116", "Timestamp Dependence", "Low", GREEN, "Static + ML"),
        ("GAS-*", "Gas Optimisation (storage in loops, etc.)", "Info", BLUE, "Gas Analyzer"),
        ("TAINT-*", "Tainted Ether Transfer / Delegatecall / Storage", "Critical", RED, "Taint"),
    ]

    headers = [
        (Inches(0.8), Inches(1.5), "SWC ID"),
        (Inches(2.3), Inches(1.5), "Vulnerability"),
        (Inches(7.0), Inches(1.5), "Severity"),
        (Inches(8.8), Inches(1.5), "Detection Layer"),
    ]
    for (left, top, text) in headers:
        add_text_box(slide, left, top, Inches(2.5), Inches(0.4),
                     text, font_size=12, color=ACCENT, bold=True)

    for i, (swc, name, severity, sev_color, layer) in enumerate(vulns):
        top = Inches(2.0) + Inches(i * 0.4)
        add_text_box(slide, Inches(0.8), top, Inches(1.3), Inches(0.35),
                     swc, font_size=11, color=GRAY)
        add_text_box(slide, Inches(2.3), top, Inches(4.5), Inches(0.35),
                     name, font_size=11, color=WHITE)
        add_text_box(slide, Inches(7.0), top, Inches(1.5), Inches(0.35),
                     severity, font_size=11, color=sev_color, bold=True)
        add_text_box(slide, Inches(8.8), top, Inches(4), Inches(0.35),
                     layer, font_size=11, color=BLUE)

    add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.5),
                 "15 static rules + CFG path analysis + taint data-flow + 10 ML classes = most comprehensive open-source Solidity auditor",
                 font_size=13, color=YELLOW, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 10 — Technology Stack
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Technology Stack")

    left_items = [
        ("Backend", "Python 3.11, FastAPI, Uvicorn, Pydantic", PURPLE),
        ("Frontend", "React 19, Axios, React-Router, Syntax Highlighter", BLUE),
        ("ML / AI", "NumPy, scikit-learn, Joblib (model serialisation)", GREEN),
        ("Analysis", "Custom Rule Engine, CFG Builder, Taint Tracker", ORANGE),
    ]
    right_items = [
        ("Testing", "pytest, FastAPI TestClient, fixture-based suites", YELLOW),
        ("Reports", "Jinja2 HTML templates, JSON, plaintext summaries", RED),
        ("Containerisation", "Docker, Docker Compose (multi-service)", GRAY),
        ("Database", "In-memory store (prod: MongoDB + Motor)", ACCENT),
    ]

    for i, (label, tech, color) in enumerate(left_items):
        top = Inches(1.8) + Inches(i * 1.1)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(5.5), Inches(0.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = color
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{label}: "
        r.font.size = Pt(15)
        r.font.color.rgb = color
        r.font.bold = True
        r2 = p.add_run()
        r2.text = tech
        r2.font.size = Pt(14)
        r2.font.color.rgb = WHITE

    for i, (label, tech, color) in enumerate(right_items):
        top = Inches(1.8) + Inches(i * 1.1)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), top, Inches(5.5), Inches(0.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = color
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{label}: "
        r.font.size = Pt(15)
        r.font.color.rgb = color
        r.font.bold = True
        r2 = p.add_run()
        r2.text = tech
        r2.font.size = Pt(14)
        r2.font.color.rgb = WHITE

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 11 — Testing & Quality Assurance
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Testing & Quality Assurance")

    test_modules = [
        ("test_static_analyzer.py", "10 tests", "Detects reentrancy, tx.origin, selfdestruct, overflow, randomness;\nverifies secure contracts produce fewer findings; field validation", GREEN),
        ("test_ml_ranker.py", "6 tests", "Predictions in [0,1] range; features are numeric;\nsecure contracts score lower than vulnerable ones", YELLOW),
        ("test_explainability.py", "3 tests", "All findings get explanations with required fields;\nedge case for empty input", ACCENT),
        ("test_api.py", "6 tests", "Health endpoint, JSON analysis, form-data analysis,\nbatch uploads, report retrieval, missing-input handling", PURPLE),
        ("test_gas_analyzer.py", "5 tests", "Gas optimisation detection; efficiency score range;\nfinding field completeness; empty-code edge case", BLUE),
        ("test_edge_cases.py", "6+ tests", "Multiline contracts, Solidity \u22650.8 skips overflow,\nstate-update-before-call not flagged as reentrancy", RED),
        ("test_store.py", "4 tests", "Save, retrieve, list, delete operations on analysis store", GRAY),
        ("conftest.py", "Fixtures", "Shared VulnerableBank, SecureBank, GasWaster contract fixtures\nreused across all test modules", ORANGE),
    ]

    for i, (name, count, desc, color) in enumerate(test_modules):
        top = Inches(1.65) + Inches(i * 0.67)
        add_text_box(slide, Inches(0.8), top, Inches(2.5), Inches(0.35),
                     name, font_size=12, color=color, bold=True)
        add_text_box(slide, Inches(3.4), top, Inches(1), Inches(0.35),
                     count, font_size=12, color=WHITE, bold=True)
        add_text_box(slide, Inches(4.5), top, Inches(8), Inches(0.6),
                     desc, font_size=10, color=GRAY)

    add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.5),
                 "40+ unit & integration tests  |  Fixture-based test architecture  |  pytest framework",
                 font_size=14, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 12 — Explainability & Reporting
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Explainability & Actionable Reporting")

    items = [
        ("Attack Scenario Narratives",
         "Every finding includes a step-by-step attack scenario explaining how an adversary "
         "could exploit the vulnerability in a real-world setting.",
         RED),
        ("Evidence & Code Context",
         "Pinpoints the exact line, matched pattern, surrounding code, and affected functions. "
         "Developers see precisely where the issue lives.",
         ORANGE),
        ("Fix Examples (Before \u2192 After)",
         "Side-by-side vulnerable vs. secure code snippets showing the recommended remediation "
         "pattern (e.g., Checks-Effects-Interactions for reentrancy).",
         GREEN),
        ("SWC Registry & Best Practice References",
         "Each finding links to the official SWC Registry entry and references OpenZeppelin "
         "libraries for battle-tested mitigations.",
         BLUE),
        ("Multi-Format Export",
         "HTML reports with severity dashboards, JSON for CI/CD integration, "
         "plain-text summaries for terminal workflows.",
         PURPLE),
    ]

    for i, (title, desc, color) in enumerate(items):
        top = Inches(1.65) + Inches(i * 1.05)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.5), Inches(0.85))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = color
        shape.line.width = Pt(1)
        add_text_box(slide, Inches(1.2), top + Inches(0.05), Inches(10.5), Inches(0.3),
                     title, font_size=14, color=color, bold=True)
        add_text_box(slide, Inches(1.2), top + Inches(0.35), Inches(10.5), Inches(0.5),
                     desc, font_size=12, color=GRAY)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 13 — Live Demo Walkthrough
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Live Demo Walkthrough")

    demos = [
        "1.  Upload VulnerableBank.sol via the React dashboard",
        "2.  Observe multi-layer scanning: Static \u2192 CFG \u2192 Taint \u2192 ML",
        "3.  View risk dashboard with severity breakdown (Critical / High / Medium / Low / Info)",
        "4.  Expand each finding: attack scenario, evidence, fix example, SWC references",
        "5.  Compare: Upload SecureBank.sol \u2014 shows significantly fewer or zero findings",
        "6.  CLI demo: python auditor.py VulnerableToken.sol -f html -o report.html",
        "7.  API demo: POST /analyze/json with curl \u2014 inspect JSON response structure",
        "8.  Run pytest test suite \u2014 show all 40+ tests passing",
    ]
    add_bullet_slide(slide, demos, start_top=Inches(1.8), font_size=15)

    add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
                 "Backend:  uvicorn app.main:app --reload    |    Frontend:  npm start    |    Tests:  pytest -v",
                 font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 14 — Project Progress & Completion
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Project Completion Status")

    progress = [
        ("Static Analysis Module (15 SWC rules)", "100%", GREEN),
        ("ML Feature Extraction & Classification (20 features, 10 classes)", "100%", GREEN),
        ("Explainability Engine (attack scenarios, fix examples)", "100%", GREEN),
        ("FastAPI REST API (6 endpoints, CORS-enabled)", "100%", GREEN),
        ("React Frontend (dark-theme dashboard, expandable findings)", "100%", GREEN),
        ("CLI Tool (single + batch, multi-format output)", "100%", GREEN),
        ("Control-Flow Graph Analyzer (path-sensitive detection)", "100%", GREEN),
        ("Taint Propagation Engine (source-to-sink tracking)", "100%", GREEN),
        ("Gas Optimisation Analyzer (efficiency scoring)", "100%", GREEN),
        ("Comprehensive Test Suite (40+ tests, 8 modules)", "100%", GREEN),
        ("Docker Containerisation (multi-service compose)", "100%", GREEN),
        ("HTML / JSON / Text Report Generation", "100%", GREEN),
    ]

    for i, (task, pct, color) in enumerate(progress):
        top = Inches(1.65) + Inches(i * 0.43)
        add_text_box(slide, Inches(0.8), top, Inches(8.5), Inches(0.4),
                     task, font_size=13, color=WHITE)

        # Progress bar background
        bar_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), top + Inches(0.08),
            Inches(2.5), Inches(0.22))
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = BG_CARD
        bar_bg.line.fill.background()

        # Progress bar fill
        bar_fill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), top + Inches(0.08),
            Inches(2.5), Inches(0.22))
        bar_fill.fill.solid()
        bar_fill.fill.fore_color.rgb = color
        bar_fill.line.fill.background()

        add_text_box(slide, Inches(12.1), top, Inches(1), Inches(0.4),
                     pct, font_size=12, color=color, bold=True)

    add_text_box(slide, Inches(0.5), Inches(6.9), Inches(12), Inches(0.5),
                 "All Planned Modules Completed \u2014 100% Project Delivery",
                 font_size=18, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 15 — Key Contributions & Novelty
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Key Contributions & Technical Novelty")

    contributions = [
        ("Multi-Layer Hybrid Analysis",
         "Unlike tools that rely on a single technique, our auditor orchestrates four independent "
         "analysis layers (static rules, CFG, taint, ML) and merges their outputs via a weighted "
         "confidence fusion algorithm \u2014 catching vulnerabilities that any single layer would miss.",
         ACCENT),
        ("Path-Sensitive + Data-Flow Aware Detection",
         "The CFG analyzer performs path-sensitive checks (e.g., external call before state update "
         "on the same execution path), while the taint engine tracks untrusted data across "
         "assignments and function boundaries \u2014 going beyond simple regex matching.",
         GREEN),
        ("Explainable AI for Smart Contract Security",
         "Every finding is enriched with a human-readable attack narrative, before/after code "
         "snippets, and authoritative SWC references \u2014 making the tool accessible to developers "
         "who are not security experts.",
         YELLOW),
        ("Full-Stack Open Auditing Platform",
         "End-to-end solution spanning a REST API, web dashboard, and CLI \u2014 enabling integration "
         "into CI/CD pipelines, one-click browser audits, and automated batch scanning.",
         PURPLE),
    ]

    for i, (title, desc, color) in enumerate(contributions):
        top = Inches(1.65) + Inches(i * 1.3)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.5), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = color
        shape.line.width = Pt(2)
        add_text_box(slide, Inches(1.2), top + Inches(0.06), Inches(10.5), Inches(0.3),
                     title, font_size=15, color=color, bold=True)
        add_text_box(slide, Inches(1.2), top + Inches(0.35), Inches(10.5), Inches(0.7),
                     desc, font_size=12, color=WHITE)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 16 — Future Scope
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_section_title(slide, "Future Scope & Enhancements")

    future = [
        "Integrate symbolic execution engines (Mythril / Manticore) for formal verification of critical paths",
        "Train deep-learning models (GNN on AST graphs) on the SmartBugs benchmark with cross-validation",
        "Add SHAP / LIME for post-hoc ML model explainability and feature attribution",
        "PDF report export with professional layout and executive summary",
        "MongoDB persistence layer for long-term analysis history and trend tracking",
        "GitHub repository URL scanning \u2014 clone, extract .sol files, audit in one click",
        "Quantitative benchmarking: precision / recall / F1 against Slither and Mythril on labelled datasets",
        "VS Code extension for inline vulnerability warnings during development",
    ]
    add_bullet_slide(slide, future, start_top=Inches(1.8), font_size=14)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 17 — Conclusion & Thank You
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.0),
                 "Thank You",
                 font_size=46, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.8),
                 "Smart Contract Auditor \u2014 Multi-Layer AI Security for Blockchain",
                 font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)

    highlights = [
        "Four-layer hybrid analysis: Static Rules + CFG + Taint Propagation + ML Classification",
        "15 SWC detectors, 10 ML vulnerability classes, 20 extracted features",
        "Path-sensitive & data-flow-aware detection beyond simple pattern matching",
        "Explainable findings with attack narratives, fix examples, and SWC references",
        "Full-stack platform: FastAPI + React + CLI + Docker",
        "40+ automated tests across 8 test modules",
    ]

    for i, line in enumerate(highlights):
        top = Inches(3.4) + Inches(i * 0.45)
        txBox = slide.shapes.add_textbox(Inches(2), top, Inches(9.5), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = "\u25c6  "
        run1.font.size = Pt(15)
        run1.font.color.rgb = ACCENT
        run1.font.bold = True
        run2 = p.add_run()
        run2.text = line
        run2.font.size = Pt(15)
        run2.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.LEFT

    add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                 "Questions?",
                 font_size=22, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

    # ── Save ─────────────────────────────────────────────────────
    output_path = os.path.join(
        os.path.dirname(__file__), "SDP_Review2_Presentation.pptx")
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    path = create_presentation()
    print(f"\nDone! Open the file: {path}")
