"""
Generate Review 1 PowerPoint Presentation for Smart Contract Auditor project.
Run: python generate_ppt.py
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme
BG_DARK = RGBColor(0x0F, 0x0F, 0x23)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0xD4, 0xFF)
WHITE = RGBColor(0xE0, 0xE0, 0xE0)
GRAY = RGBColor(0xA0, 0xA0, 0xB0)
RED = RGBColor(0xFF, 0x47, 0x57)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
YELLOW = RGBColor(0xFF, 0xA5, 0x02)
GREEN = RGBColor(0x2E, 0xD5, 0x73)
BLUE = RGBColor(0x1E, 0x90, 0xFF)
PURPLE = RGBColor(0x7C, 0x5C, 0xFF)


def set_slide_bg(slide, color=BG_DARK):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_bullet_slide(slide, items, start_top=Inches(1.8), font_size=16, color=WHITE, bullet_color=ACCENT):
    for i, item in enumerate(items):
        top = start_top + Inches(i * 0.5)
        txBox = slide.shapes.add_textbox(Inches(1), top, Inches(8), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        # Bullet point
        run1 = p.add_run()
        run1.text = "  >  "
        run1.font.size = Pt(font_size)
        run1.font.color.rgb = bullet_color
        run1.font.bold = True

        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = color


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ========== SLIDE 1: Title Slide ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide)

    # Title
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                 "Smart Contract Auditor", font_size=44, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
                 "AI-Powered Security Analysis for Solidity Smart Contracts",
                 font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
                 "Using Blockchain + Machine Learning + Static Analysis",
                 font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Subtitle info
    add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
                 "Software Design Project - Review 1",
                 font_size=18, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
                 "Student: Harish  |  Guide: [Your Guide Name]",
                 font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
                 "Department of Computer Science and Engineering",
                 font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ========== SLIDE 2: Problem Statement ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "Problem Statement", font_size=34, color=ACCENT, bold=True)

    # Divider line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(3), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    items = [
        "Smart contract bugs cause catastrophic financial losses (DAO hack: $60M lost)",
        "DeFi protocols hold billions in TVL - security is critical",
        "Existing tools (Slither, Mythril) produce many false positives",
        "Lack of human-friendly remediation guidance in current tools",
        "Manual audits are expensive ($5K-$50K+) and time-consuming",
        "Need: AI-powered auditor that reduces false positives & explains findings",
    ]
    add_bullet_slide(slide, items, start_top=Inches(1.8))

    # Stats box
    add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(1),
                 "$3.8 Billion lost to smart contract exploits in 2022 alone (Chainalysis Report)",
                 font_size=16, color=RED, bold=True, alignment=PP_ALIGN.CENTER)

    # ========== SLIDE 3: Objectives ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "Project Objectives", font_size=34, color=ACCENT, bold=True)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(3), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    items = [
        "Build a pipeline to ingest Solidity code, parse AST, and produce analyzable representations",
        "Implement a static analysis module with pattern-based detection for SWC vulnerabilities",
        "Integrate ML classifiers to reduce false positives and predict vulnerability severity",
        "Provide explainability (code highlights, attack scenarios, fix examples) for each finding",
        "Generate actionable remediation suggestions with code snippets",
        "Expose the auditor via REST API and modern web UI for uploads and reports",
        "Build CLI tool for automated batch scans and CI/CD integration",
    ]
    add_bullet_slide(slide, items, start_top=Inches(1.8), font_size=15)

    # ========== SLIDE 4: Architecture ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "System Architecture", font_size=34, color=ACCENT, bold=True)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(3), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    # Architecture components as boxes
    components = [
        (Inches(0.5), Inches(2.0), "React Frontend\n(Web UI)", BLUE),
        (Inches(3.5), Inches(2.0), "FastAPI Backend\n(REST API)", PURPLE),
        (Inches(6.5), Inches(2.0), "Static Analyzer\n(Pattern Rules)", ORANGE),
        (Inches(9.5), Inches(2.0), "ML Classifier\n(Risk Scoring)", GREEN),
        (Inches(3.5), Inches(4.0), "Explainability\nEngine", YELLOW),
        (Inches(6.5), Inches(4.0), "Report Generator\n(HTML/JSON/PDF)", RED),
        (Inches(0.5), Inches(4.0), "CLI Tool\n(Batch Scanning)", GRAY),
        (Inches(9.5), Inches(4.0), "Analysis Store\n(MongoDB/Memory)", ACCENT),
    ]

    for (left, top, text, color) in components:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(1.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        tf = shape.text_frame
        tf.word_wrap = True
        for i, line in enumerate(text.split('\n')):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(13) if i == 0 else Pt(11)
            p.font.color.rgb = color if i == 0 else GRAY
            p.font.bold = (i == 0)
            p.alignment = PP_ALIGN.CENTER

    # Flow arrows description
    add_text_box(slide, Inches(0.5), Inches(5.8), Inches(12), Inches(1),
                 "Flow: Upload .sol -> Static Analysis + ML Classification -> Merge & Rank -> Explain -> Generate Report",
                 font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ========== SLIDE 5: Tech Stack ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "Technology Stack", font_size=34, color=ACCENT, bold=True)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(3), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    # Two-column layout
    left_items = [
        ("Backend", "Python, FastAPI, Uvicorn", PURPLE),
        ("Frontend", "React.js, CSS3", BLUE),
        ("ML/AI", "NumPy, scikit-learn", GREEN),
        ("Analysis", "Custom Rule Engine, Regex Patterns", ORANGE),
    ]

    right_items = [
        ("Database", "In-memory Store / MongoDB", ACCENT),
        ("CLI", "Python argparse, Batch processing", GRAY),
        ("Reports", "HTML, JSON, Text formats", YELLOW),
        ("Container", "Docker, Docker Compose", RED),
    ]

    for i, (label, tech, color) in enumerate(left_items):
        top = Inches(1.8) + Inches(i * 1.1)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(5.5), Inches(0.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = color
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{label}: "
        run.font.size = Pt(15)
        run.font.color.rgb = color
        run.font.bold = True
        run2 = p.add_run()
        run2.text = tech
        run2.font.size = Pt(14)
        run2.font.color.rgb = WHITE

    for i, (label, tech, color) in enumerate(right_items):
        top = Inches(1.8) + Inches(i * 1.1)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), top, Inches(5.5), Inches(0.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = color
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{label}: "
        run.font.size = Pt(15)
        run.font.color.rgb = color
        run.font.bold = True
        run2 = p.add_run()
        run2.text = tech
        run2.font.size = Pt(14)
        run2.font.color.rgb = WHITE

    # ========== SLIDE 6: Modules Implemented ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "Modules Implemented (Review 1)", font_size=34, color=ACCENT, bold=True)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(3), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    modules = [
        ("Static Analyzer", "15 rule-based detectors covering SWC Registry vulnerabilities\n(Reentrancy, Integer Overflow, tx.origin, Unchecked Returns, etc.)", GREEN, "DONE"),
        ("ML Classifier", "Feature extraction from Solidity code (20 features)\nHeuristic-based scoring across 10 vulnerability classes", GREEN, "DONE"),
        ("Explainability Engine", "Attack scenarios, fix examples, code context highlighting\nReferences to SWC Registry and best practices", GREEN, "DONE"),
        ("Report Generator", "HTML reports with severity stats, detailed findings\nJSON and text summary formats", GREEN, "DONE"),
        ("FastAPI Backend", "REST API with /analyze, /report, /batch endpoints\nCORS-enabled for React frontend", GREEN, "DONE"),
        ("React Frontend", "Modern dark-theme UI with code editor, sample contracts\nReal-time analysis results with expandable explanations", GREEN, "DONE"),
        ("CLI Tool", "Command-line scanner with text/JSON/HTML output\nBatch analysis support", GREEN, "DONE"),
    ]

    for i, (name, desc, color, status) in enumerate(modules):
        top = Inches(1.7) + Inches(i * 0.75)
        # Status indicator
        add_text_box(slide, Inches(0.8), top, Inches(0.7), Inches(0.5),
                     status, font_size=10, color=GREEN, bold=True)

        add_text_box(slide, Inches(1.5), top, Inches(2.5), Inches(0.5),
                     name, font_size=15, color=ACCENT, bold=True)

        add_text_box(slide, Inches(4.2), top, Inches(8), Inches(0.7),
                     desc, font_size=11, color=GRAY)

    # ========== SLIDE 7: Vulnerability Detection ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "Vulnerability Detection Capabilities", font_size=34, color=ACCENT, bold=True)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(3), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    vulns = [
        ("SWC-107", "Reentrancy", "Critical", RED),
        ("SWC-101", "Integer Overflow/Underflow", "High", ORANGE),
        ("SWC-115", "tx.origin Authorization", "High", ORANGE),
        ("SWC-104", "Unchecked Call Return Value", "High", ORANGE),
        ("SWC-105", "Unprotected Ether Withdrawal", "Critical", RED),
        ("SWC-106", "Unprotected SELFDESTRUCT", "Critical", RED),
        ("SWC-112", "Delegatecall to Untrusted Callee", "Critical", RED),
        ("SWC-120", "Weak Sources of Randomness", "High", ORANGE),
        ("SWC-103", "Floating Pragma", "Info", BLUE),
        ("SWC-111", "Deprecated Functions", "Low", GREEN),
    ]

    # Table-like layout
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(1.5), Inches(0.4),
                 "SWC ID", font_size=12, color=ACCENT, bold=True)
    add_text_box(slide, Inches(2.5), Inches(1.6), Inches(4), Inches(0.4),
                 "Vulnerability", font_size=12, color=ACCENT, bold=True)
    add_text_box(slide, Inches(7), Inches(1.6), Inches(2), Inches(0.4),
                 "Severity", font_size=12, color=ACCENT, bold=True)

    for i, (swc_id, name, severity, color) in enumerate(vulns):
        top = Inches(2.1) + Inches(i * 0.45)
        add_text_box(slide, Inches(0.8), top, Inches(1.5), Inches(0.4),
                     swc_id, font_size=12, color=GRAY)
        add_text_box(slide, Inches(2.5), top, Inches(4), Inches(0.4),
                     name, font_size=12, color=WHITE)
        add_text_box(slide, Inches(7), top, Inches(2), Inches(0.4),
                     severity, font_size=12, color=color, bold=True)

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
                 "+ ML-based detection for: Front Running, Denial of Service, Access Control patterns",
                 font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ========== SLIDE 8: ML Pipeline ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "ML Classification Pipeline", font_size=34, color=ACCENT, bold=True)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(3), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    # Pipeline steps
    steps = [
        ("1. Feature Extraction", "20 features extracted from Solidity source code:\n- Code complexity (LOC, function count, modifier count)\n- External interactions (call, transfer, send, delegatecall counts)\n- Security patterns (ReentrancyGuard, Ownable, SafeMath)\n- Solidity version, pragma type, state variables, loops"),
        ("2. Vulnerability Classification", "10 vulnerability classes predicted:\nReentrancy, Integer Overflow, Unchecked Return,\nAccess Control, Timestamp Dependence, tx.origin,\nSELFDESTRUCT, Delegatecall, Front Running, DoS"),
        ("3. Hybrid Scoring", "Combined confidence = 60% Static + 40% ML\nML-only findings added when confidence > 0.6\nResults sorted by confidence descending"),
    ]

    for i, (title, desc) in enumerate(steps):
        top = Inches(1.8) + Inches(i * 1.7)

        # Step box
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11), Inches(1.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = ACCENT
        shape.line.width = Pt(1)

        add_text_box(slide, Inches(1.2), top + Inches(0.1), Inches(10), Inches(0.4),
                     title, font_size=16, color=ACCENT, bold=True)
        add_text_box(slide, Inches(1.2), top + Inches(0.45), Inches(10), Inches(0.9),
                     desc, font_size=12, color=GRAY)

    # ========== SLIDE 9: Demo Screenshots Description ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "Live Demo", font_size=34, color=ACCENT, bold=True)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(3), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    demos = [
        "1. Upload a vulnerable Solidity contract (VulnerableBank.sol)",
        "2. Run AI-powered analysis - see real-time scanning progress",
        "3. View risk dashboard with severity breakdown",
        "4. Explore each vulnerability with detailed explanations",
        "5. See attack scenarios and fix examples for each finding",
        "6. Compare: Analyze SecureBank.sol - shows fewer/no vulnerabilities",
        "7. CLI Demo: Run batch analysis from command line",
        "8. API Demo: Show REST API endpoints and responses",
    ]
    add_bullet_slide(slide, demos, start_top=Inches(1.8), font_size=16)

    add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
                 "Backend: python -m uvicorn app.main:app --reload    |    Frontend: npm start",
                 font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ========== SLIDE 10: Project Progress ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "Project Progress & Timeline", font_size=34, color=ACCENT, bold=True)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(3), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    # Progress items with status
    progress = [
        ("Requirements & Environment Setup", "100%", GREEN),
        ("Ingest Pipeline & AST Extraction", "100%", GREEN),
        ("Static Analysis Module (15 rules)", "100%", GREEN),
        ("ML Feature Extraction & Classification", "100%", GREEN),
        ("Explainability Engine", "100%", GREEN),
        ("FastAPI Backend + REST API", "100%", GREEN),
        ("React Frontend UI", "100%", GREEN),
        ("CLI Tool", "100%", GREEN),
        ("Symbolic Execution (Mythril integration)", "0%", RED),
        ("Advanced ML (Deep Learning, SHAP/LIME)", "0%", ORANGE),
        ("PDF Report Export", "0%", ORANGE),
        ("Evaluation & Benchmarking", "0%", ORANGE),
    ]

    for i, (task, pct, color) in enumerate(progress):
        top = Inches(1.7) + Inches(i * 0.43)
        add_text_box(slide, Inches(0.8), top, Inches(7), Inches(0.4),
                     task, font_size=13, color=WHITE)
        add_text_box(slide, Inches(8.5), top, Inches(1.5), Inches(0.4),
                     pct, font_size=13, color=color, bold=True)

    # Overall progress bar label
    add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5),
                 "Overall Progress: ~65% Complete for Review 1",
                 font_size=16, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

    # ========== SLIDE 11: Remaining Work ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                 "Remaining Work (Review 2)", font_size=34, color=ACCENT, bold=True)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(3), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    remaining = [
        "Integrate Mythril/Manticore for symbolic execution & deeper analysis",
        "Train proper ML models on SmartBugs dataset with cross-validation",
        "Implement SHAP/LIME for ML model explainability",
        "Add PDF report export with professional formatting",
        "Implement MongoDB persistence for analysis history",
        "Add GitHub repo URL scanning support",
        "Evaluation: Compare precision/recall against Slither and Mythril baselines",
        "Docker containerization for easy deployment",
        "User study with 3-5 developers for qualitative feedback",
    ]
    add_bullet_slide(slide, remaining, start_top=Inches(1.8), font_size=15)

    # ========== SLIDE 12: Conclusion ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                 "Thank You", font_size=44, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
                 "Smart Contract Auditor - AI-Powered Security for Blockchain",
                 font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)

    items_text = (
        "Hybrid Analysis: Static Rules + ML Classification\n"
        "15 SWC vulnerability detectors with explainable findings\n"
        "Full-stack implementation: FastAPI + React + CLI\n"
        "Ready for live demo"
    )

    txBox = slide.shapes.add_textbox(Inches(2), Inches(4.0), Inches(9), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in items_text.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)

    add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                 "Questions?",
                 font_size=20, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

    # Save
    output_path = os.path.join(os.path.dirname(__file__), "..", "docs", "SDP_Review1_Presentation.pptx")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    path = create_presentation()
    print(f"\nDone! Open the file: {path}")
